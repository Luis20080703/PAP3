from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Definir cores
COR_AZUL = RGBColor(37, 99, 235)  # Blue-600
COR_TEXTO = RGBColor(17, 24, 39)  # Gray-900
COR_SUBTITULO = RGBColor(75, 85, 99)  # Gray-600
COR_FUNDO_CLARO = RGBColor(243, 244, 246)  # Gray-100

def adicionar_titulo_slide(prs, titulo, subtitulo=""):
    """Cria um slide com título e subtítulo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Fundo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Barra azul no topo
    barra = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.5))
    barra.fill.solid()
    barra.fill.fore_color.rgb = COR_AZUL
    barra.line.color.rgb = COR_AZUL
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    
    # Subtítulo
    if subtitulo:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitulo
        p2.font.size = Pt(32)
        p2.font.color.rgb = COR_SUBTITULO
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def adicionar_conteudo_slide(prs, titulo, conteudo_items):
    """Cria um slide com título e pontos de conteúdo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fundo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Barra azul no topo
    barra = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    barra.fill.solid()
    barra.fill.fore_color.rgb = COR_AZUL
    barra.line.color.rgb = COR_AZUL
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Conteúdo
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, item in enumerate(conteudo_items):
        if i > 0:
            tf2.add_paragraph()
        p = tf2.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COR_TEXTO
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

# SLIDE 1: Capa
slide1 = adicionar_titulo_slide(prs, "🎯 Rotas no NexusHand", "Como funciona a navegação na aplicação React")

# SLIDE 2: Introdução
slide2 = adicionar_conteudo_slide(prs, "Introdução", [
    "📱 NexusHand é uma plataforma para atletas e treinadores de andebol",
    "🎬 Partilha de jogadas, dicas técnicas e estatísticas",
    "🔐 Sistema de autenticação com 3 páginas principais",
    "⚡ Usa Estado & Contexto (sem React Router)"
])

# SLIDE 3: Estrutura
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

barra = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
barra.fill.solid()
barra.fill.fore_color.rgb = COR_AZUL
barra.line.color.rgb = COR_AZUL

txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "As 3 Rotas Principais"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Três caixas
rotas = [
    ("📄 HOME", "Página Inicial", "• Landing page\n• Botão 'Entrar'\n• Sem autenticação"),
    ("🔐 LOGIN", "Autenticação", "• Formulário\n• Login/Registo\n• Sem autenticação"),
    ("📊 DASHBOARD", "Painel Principal", "• Jogadas\n• Dicas\n• Estatísticas\n• Protegida")
]

x_positions = [0.5, 3.5, 6.5]
for idx, (titulo, subtitulo, desc) in enumerate(rotas):
    # Box
    box = slide3.shapes.add_shape(1, Inches(x_positions[idx]), Inches(1.2), Inches(2.8), Inches(5.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COR_FUNDO_CLARO
    box.line.color.rgb = COR_AZUL
    box.line.width = Pt(2)
    
    # Título
    txBox_t = slide3.shapes.add_textbox(Inches(x_positions[idx] + 0.1), Inches(1.4), Inches(2.6), Inches(0.8))
    tf_t = txBox_t.text_frame
    p_t = tf_t.paragraphs[0]
    p_t.text = titulo
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = COR_AZUL
    p_t.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    txBox_s = slide3.shapes.add_textbox(Inches(x_positions[idx] + 0.1), Inches(2.2), Inches(2.6), Inches(0.5))
    tf_s = txBox_s.text_frame
    p_s = tf_s.paragraphs[0]
    p_s.text = subtitulo
    p_s.font.size = Pt(14)
    p_s.font.italic = True
    p_s.font.color.rgb = COR_SUBTITULO
    p_s.alignment = PP_ALIGN.CENTER
    
    # Descrição
    txBox_d = slide3.shapes.add_textbox(Inches(x_positions[idx] + 0.15), Inches(2.9), Inches(2.5), Inches(3.5))
    tf_d = txBox_d.text_frame
    tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]
    p_d.text = desc
    p_d.font.size = Pt(13)
    p_d.font.color.rgb = COR_TEXTO
    p_d.alignment = PP_ALIGN.LEFT

# SLIDE 4: Fluxo de Navegação
slide4 = adicionar_conteudo_slide(prs, "Fluxo de Navegação", [
    "1️⃣ Aplicação inicia → Home (user = null)",
    "2️⃣ Utilizador clica 'Entrar' → Login",
    "3️⃣ Faz login com sucesso → user é preenchido",
    "4️⃣ Automaticamente vai para → Dashboard",
    "5️⃣ Clica 'Sair' → Logout",
    "6️⃣ user é zerado → Volta para Home"
])

# SLIDE 5: Como Funciona o Sistema
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

barra = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
barra.fill.solid()
barra.fill.fore_color.rgb = COR_AZUL
barra.line.color.rgb = COR_AZUL

txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Como Funciona?"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Código
codigo = """const [currentPage, setCurrentPage] = useState('home');
const { user } = useApp();

return (
  <>
    {currentPage === 'home' && !user && 
      <Home onNavigateToLogin={() => setCurrentPage('login')} />
    }
    
    {currentPage === 'login' && !user && 
      <Login onBack={() => setCurrentPage('home')} />
    }
    
    {user && 
      <Dashboard onLogout={handleLogout} />
    }
  </>
);"""

txBox_code = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
tf_code = txBox_code.text_frame
tf_code.word_wrap = True
p_code = tf_code.paragraphs[0]
p_code.text = codigo
p_code.font.size = Pt(11)
p_code.font.name = "Courier New"
p_code.font.color.rgb = COR_TEXTO
tf_code.margin_bottom = Inches(0.1)
tf_code.margin_top = Inches(0.1)
tf_code.margin_left = Inches(0.2)
tf_code.margin_right = Inches(0.2)

# SLIDE 6: Estado & Contexto
slide6 = adicionar_conteudo_slide(prs, "Estado & Contexto", [
    "🎛️ ESTADO LOCAL (currentPage):",
    "   • Controla qual página mostrar",
    "   • Valores: 'home' | 'login' | 'dashboard'",
    "",
    "🔐 CONTEXTO GLOBAL (user):",
    "   • Dados do utilizador autenticado",
    "   • null se não autenticado",
    "   • Sincroniza automaticamente com a navegação"
])

# SLIDE 7: Renderização Condicional
slide7 = adicionar_conteudo_slide(prs, "Renderização Condicional", [
    "✅ SE currentPage === 'home' E !user",
    "   → Mostra <Home />",
    "",
    "✅ SE currentPage === 'login' E !user",
    "   → Mostra <Login />",
    "",
    "✅ SE user existe",
    "   → Mostra <Dashboard />"
])

# SLIDE 8: Vantagens
slide8 = adicionar_conteudo_slide(prs, "✅ Vantagens", [
    "🎯 Simples → Sem dependências externas",
    "⚡ Rápido → Re-renderiza apenas o necessário",
    "🧠 Fácil de entender → Lógica clara e direta",
    "🔄 Automático → Sincronização do user/navegação",
    "📱 Lightweight → Menos código, menos peso"
])

# SLIDE 9: Limitações
slide9 = adicionar_conteudo_slide(prs, "❌ Limitações", [
    "🔗 Sem URLs reais → /home, /login, /dashboard não existem",
    "⏮️ Sem histórico → Botão 'voltar' não funciona",
    "🔖 Sem bookmarking → Não podes partilhar URLs",
    "🔍 Difícil SEO → Motores de busca não indexam",
    "📈 Não escala → Limitado para muitas rotas"
])

# SLIDE 10: Solução: React Router
slide10 = adicionar_conteudo_slide(prs, "Alternativa: React Router", [
    "Se quisermos URLs reais e histórico de browser:",
    "",
    "🔧 npm install react-router-dom",
    "",
    "Benefícios:",
    "  ✅ URLs reais (/home, /login, /dashboard)",
    "  ✅ Histórico de browser (botão voltar)",
    "  ✅ Bookmarking de URLs",
    "  ✅ Melhor para SEO"
])

# SLIDE 11: Ficheiros Envolvidos
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide11.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

barra = slide11.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
barra.fill.solid()
barra.fill.fore_color.rgb = COR_AZUL
barra.line.color.rgb = COR_AZUL

txBox = slide11.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Ficheiros Envolvidos"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

estrutura = """src/
├── App.tsx ......................... 🎛️ Centro de controlo
├── context/AppContext.tsx .......... 🔐 Autenticação & user
├── components/
│   ├── Home.tsx ................... 📄 Rota 1
│   ├── Login.tsx .................. 🔐 Rota 2
│   └── Dashboard.tsx .............. 📊 Rota 3
│   ├── PlaysSection.tsx
│   ├── TipsSection.tsx
│   ├── TeamStatsSection.tsx
│   └── AthleteStatsSection.tsx
└── services/api.ts ................ 📡 Chamadas à API"""

txBox_est = slide11.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
tf_est = txBox_est.text_frame
tf_est.word_wrap = True
p_est = tf_est.paragraphs[0]
p_est.text = estrutura
p_est.font.size = Pt(12)
p_est.font.name = "Courier New"
p_est.font.color.rgb = COR_TEXTO

# SLIDE 12: Exemplo Prático
slide12 = adicionar_conteudo_slide(prs, "Exemplo Prático: Login", [
    "1️⃣ Utilizador clica 'Entrar' na Home",
    "   → onNavigateToLogin() é chamada",
    "   → setCurrentPage('login')",
    "",
    "2️⃣ Insere email e password, clica 'Entrar'",
    "   → login(email, password) é chamada",
    "   → Envia dados à API",
    "",
    "3️⃣ API retorna dados do utilizador",
    "   → setUser(userData) no contexto",
    "   → user agora tem dados ✅",
    "",
    "4️⃣ React detecta mudança em 'user'",
    "   → Re-renderiza App.tsx",
    "   → {user && <Dashboard />} → VERDADEIRO",
    "   → Mostra Dashboard ✅"
])

# SLIDE 13: Perguntas Frequentes
slide13 = adicionar_conteudo_slide(prs, "❓ Perguntas Frequentes", [
    "P: E se recarregar a página?",
    "R: useEffect verifica se user existe. Se sim, fica no Dashboard.",
    "",
    "P: Por que não usam React Router?",
    "R: Para simplicidade. Router é melhor para apps maiores.",
    "",
    "P: Como protegem as rotas?",
    "R: {user && <Dashboard />} - só renderiza se user existe."
])

# SLIDE 14: Conclusão
slide14 = adicionar_titulo_slide(prs, "🎓 Resumo", "O sistema de rotas é simples mas eficaz:\nEstado Local + Contexto Global + Renderização Condicional = Navegação")

# Salvar
prs.save(r'c:\PAP3\Andebolonlineplatform-main\ROTAS_APRESENTACAO.pptx')
print("✅ Apresentação criada com sucesso: ROTAS_APRESENTACAO.pptx")
